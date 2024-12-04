"""
Created on Mon Dec 2 2024

@author: LukeMurphy
"""

# python-pptx package:
# https://python-pptx.readthedocs.io/en/latest/

###############################################################################
###############################################################################
###############################################################################


####### FOLDER #######

# Base folder location
# this folder contains images that will be used in the output, and the ppt template
folder = "C:/Users/LukeMurphy/OneDrive - Dense Air/Data Science/Business Development/Automated collateral/"


###############################################################################
###############################################################################
###############################################################################

####### LIBS, FOLDERS, CREDS #######

import argparse
import json

import pandas as pd
import google.auth
import geopandas as gpd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
import cartopy.crs as ccrs
import cartopy.io.img_tiles as cimgt
import io
import math

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from google.oauth2 import service_account
from google.cloud import bigquery
from shapely import wkt
from urllib.request import urlopen, Request
from PIL import Image
from datetime import datetime
from pathlib import Path

## PowerPoint Template
template = "Building_analysis_template v4.pptx"

## GCP project JSON key file
JSON_KEY_FILE = folder + "Code/" + "YOUR_KEY_FILE_NAME.json" 

## SET UP GCP CREDENTIALS - LOAD FROM JSON FILE
credentials = service_account.Credentials.from_service_account_file(
    JSON_KEY_FILE,
    scopes = ["https://www.googleapis.com/auth/cloud-platform"]
    )

# Use the credentials to authenticate
try:
    # Explicitly refresh the credentials to ensure they're valid.
    credentials.refresh(google.auth.transport.requests.Request())
    client = bigquery.Client.from_service_account_json(JSON_KEY_FILE)
    print("1) GCP has authenticated succesfully")
except Exception as e:
    print("1) Error authenticating:\n\t",e)


###############################################################################
###############################################################################
###############################################################################


# PARSE CLI ARGUMENTS
parser = argparse.ArgumentParser()
parser.add_argument('-c', '--config-file', type=str, required=True,
                    help="Configuration file. REQUIRED.", dest='config_file')
args = parser.parse_args()

with open(args.config_file) as f:
    config = json.loads(f.read())

# set constants
    
# DW building ID(s)
BOI_list = config['building_ids']

# Calculate hexbin area from size provided in config
hexbin_diag = config['hexbin_diag']
hexbin_size = round(3*((math.sqrt(3))*hexbin_diag**2) / 8)

# Limit to the number of floors to be included in charts and maps
max_flrs = config['max_floors']
    

###############################################################################
###############################################################################
###############################################################################


####### SET OF STATIC IMAGE FILES AND COLOURS  #######

## Hexbin size graphic file
hxbin_size_image = folder + 'Images/' + 'Internal_hexbin_size_v2.png'

## Measurement scale graphic file
meas_scale_image = folder + 'Images/' + 'Measurement_colour_scale.png'

## RF scale graphic file
RF_scale_image = folder + 'Images/' + 'RF_colour_scales.png'

# DW colours for measurements
meas_bins = [-np.inf, 0, 10, 50, 100, 200, 500, 1000, 2000, 5000, 10000, np.inf]
meas_labels = ['#FFFFFF', '#FFDE00', '#FFBC00', '#FF9800', '#FF6C27', '#FF2447', '#D90065', '#AD0075', '#820080', '#510081', '#000080']

# DW colours for RSRP
rsrp_bins = [-np.inf, -124, -120, -116, -112, -108, -104, -100, -40, np.inf]
rsrp_labels = ["#381D1D","#850000","#FF0000","#FF8500","#FFFF00","#85FF85","#4ABE45","#008000","#FFFFFF"]

# DW colours for RSSNR
rssnr_bins = [-np.inf, -6, -4, -2, 0, 2, 5, 10, 30, np.inf]
rssnr_labels = ["#381D1D","#850000","#FF0000","#FF8500","#FFFF00","#85FF85","#4ABE45","#008000","#FFFFFF"]

# Additional marketing/experience colours for RSRP
mkt_rsrp_bins = [-np.inf, -110, -98, -85, np.inf]
mkt_rsrp_labels = ["#9C0000","#C01200","#FDC000","#00B050"]

# Additional marketing/experience colours for RSSNR
mkt_rssnr_bins = [-np.inf, -2, 4, 10, np.inf]
mkt_rssnr_labels = ["#9C0000","#C01200","#FDC000","#00B050"]


# 25th ptile calc
def q25(x):
    return x.quantile(0.25)
   
    

###############################################################################
###############################################################################
###############################################################################


####### GENERATE ANALYSIS  #######


for i, building_id in enumerate(BOI_list): 


    ###############################################################################
    ###############################################################################
    ###############################################################################
    
    
    ####### BUILDING DETAILS - REQUIRES INPUT OF DESIRED BUILDING ID #######
    
    SQL_building_info = """
    # AUTOMATED RETRIEVAL OF BUILDING INFORMATION
    WITH all_ids AS (
            SELECT  *
            FROM    `denseware-prod-1-d7b8.d_denseware_lead_generation_us_prod.t_costar-flat_v1`
            )
    
    # subset to select building
    ,my_build AS (
            SELECT  *
            FROM    all_ids
            WHERE   building_id = """ + str(building_id) + """
    )
    
    SELECT  my_build.building_id
            ,my_build.costar_property_address
            ,my_build.costar_city
            ,my_build.costar_state_code
            ,my_build.costar_zip
            ,my_build.costar_rba
            ,my_build.costar_year_built
            ,my_build.costar_number_of_parking_spaces
            ,my_build.costar_number_of_stories
            ,my_build.denseair_built_or_renovated
            ,alt.altitude_low_metres
            FROM    my_build AS my_build
            LEFT JOIN `denseware-prod-1-d7b8.d_denseware_reference_us_prod.t_building_z_axis` AS alt
            ON my_build.building_id = alt.building_id
    ;
    """
    
    # retrieve Big Query SQL as dataframe
    building_info = client.query(SQL_building_info).to_dataframe()
    
    building_info = building_info.head(1)
    
    ## Address
    address = building_info['costar_property_address'].item()
    city = building_info['costar_city'].item()
    state = building_info['costar_state_code'].item()
    zip_code = building_info['costar_zip'].item()
    stories = building_info['costar_number_of_stories'].item()
    rba = building_info['costar_rba'].item()
    year = building_info['denseair_built_or_renovated'].item()
    low_alt = building_info['altitude_low_metres'].item() # Z-axis derived 
    
    
    ###############################################################################
    ###############################################################################
    ###############################################################################
    
    
    ####### GENERATE REQUIRED BUILDING GRID #######
    
    
    ## Retrieve XY grid for building
    SQL_building_grid = """
    
    #############################################
    # SET PARAMETERS AND USER DEFINED FUNCTIONS
    #############################################
    
    
    DECLARE area_sqft FLOAT64;
    DECLARE area FLOAT64;     
    DECLARE area_multiplier FLOAT64;
    DECLARE R FLOAT64;              
    DECLARE step_x FLOAT64;         
    DECLARE step_y FLOAT64;         
    DECLARE offset_x FLOAT64;       
    DECLARE offset_y FLOAT64;       
    DECLARE hb_buffer FLOAT64;      
    
    
    -- MANUAL QUERY PARAMETERS
    SET area_sqft = """ + str(hexbin_size) + """;           
    SET area_multiplier = 2.0;      
    SET hb_buffer = 0;              
    
    
    -- DERIVED QUERY PARAMETERS
    SET area = area_sqft * 0.092903;
    SET R = SQRT(2 * SQRT(3) * area) / 3;
    SET step_x = 3 * R;
    SET step_y = SQRT(3) * R;
    SET offset_x = 1.5 * R;
    SET offset_y = SQRT(3) / 2 * R;
    
    
    -- UDFs
    CREATE TEMP FUNCTION deg_to_rad(a FLOAT64)
    RETURNS FLOAT64
    AS (
      a / 90 * ASIN(1)
    );
    
    
    CREATE TEMP FUNCTION rad_to_deg(a FLOAT64)
    RETURNS FLOAT64
    AS (
      a * 90 / ASIN(1)
    );
    
    
    -- this function returns an approximate radius of the earth at a given latitude
    -- a is radius at equator; b is radius at poles
    CREATE TEMP FUNCTION earth_radius(lat FLOAT64)
    RETURNS FLOAT64
    AS (
      (
        SELECT    SQRT(
                  (POWER(a * a * COS(r), 2) + POWER(b * b * SIN(r), 2))
                  /
                  (POWER(a * COS(r), 2) + POWER(b * SIN(r), 2))
                )
        FROM      (SELECT ABS(deg_to_rad(lat)) AS r, 6378137 AS a, 6356752 AS b)
      )
    );
    
    
    ###############################################
    # PREPARATION WORK FOR ALIGNING HEXBINS
    ###############################################
    
    
    -- estimate building orientation for aligning hexbins;
    -- multipolygons are excluded at this point as existing process does not cater for them
    
    CREATE TEMPORARY TABLE buildings AS (
    
    WITH BOI AS (
      SELECT  building_id
              ,boundary
              ,boundary_buffered
              ,area_m2
              ,CASE WHEN area_m2 >= area * area_multiplier THEN 0
                ELSE 1
                  END AS display_flag
      FROM    `denseware-prod-1-d7b8.d_denseware_reference_us_prod.t_denseware_building`
      WHERE building_id = """ + str(building_id) + """
      )
    
      ,points AS (
      -- get outer ring node geometries
      SELECT      building_id
                  , boundary
                  , boundary_buffered
                  , area_m2
                  , point
                  , node
      FROM        BOI
                  , UNNEST(REGEXP_EXTRACT_ALL(ST_ASTEXT(ST_EXTERIORRING(boundary)), r'[^,\(\)]+')) AS point WITH OFFSET AS node
      WHERE       1 = 1
      -- MULTIPOLYGONS won't work so are excluded - this edge case isn't managed here as it is understood
      -- Jon will address it in the source data at some point
      AND         ST_GEOMETRYTYPE(boundary) = 'ST_Polygon' 
      )
      
      , vertices AS (
      -- this gets all the outer ring vertices
      SELECT      building_id
                  , boundary
                  , boundary_buffered
                  , area_m2
                  , ST_CENTROID(boundary) AS center
                  , ST_GEOGFROMTEXT(CONCAT("POINT(", point, ")")) AS vertice
                  , node - 1 as node
      FROM        points
      WHERE       node > 0
      )
      
      , sides AS (
      -- this get the sides of the outer ring
      SELECT        building_id
                    , boundary
                    , boundary_buffered
                    , area_m2
                    , node
                    , ST_AZIMUTH(vertice, FIRST_VALUE(vertice) OVER edge) AS azimuth
                    , ST_MAKELINE(vertice, FIRST_VALUE(vertice) OVER edge) AS side
      FROM          vertices
      WHERE         1 = 1
      WINDOW        edge AS (PARTITION BY building_id ORDER BY building_id, node ROWS BETWEEN 1 PRECEDING AND CURRENT ROW)
      )
      
      , azimuths AS (
      -- the side with greatest length will be used for building azimuth
      SELECT        building_id
                    , boundary
                    , boundary_buffered
                    , area_m2
                    , node
                    , azimuth
                    , side
                    , ROW_NUMBER() OVER (PARTITION BY building_id ORDER BY ST_LENGTH(side) DESC) AS row_num
      FROM          sides
      WHERE         1 = 1
      AND           node > 0
      )
    
    -- calculate azimuth
    SELECT        building_id
                  , boundary
                  , boundary_buffered
                  , area_m2
                  , CASE 
                      WHEN ROUND(azimuth / ACOS(-1) * 180 - 180, 0) = 0 THEN 0
                      ELSE CASE
                        WHEN azimuth <= ACOS(-1) THEN ROUND(azimuth / ACOS(-1) * 180, 0)
                        ELSE ROUND(azimuth / ACOS(-1) * 180 - 180, 0)
                        END
                      END AS azimuth
    FROM          azimuths
    WHERE         1 = 1
    AND           row_num = 1
    );
    
    ###############################################
    # SPLIT BUILDINGS IN XY
    ###############################################
    
    create or replace table `da-datascience-pt.LM_sandpit.Temp_LG_building_grid` as (
    WITH bboxes AS (
      SELECT          building_id
                      , ST_CENTROID(boundary) AS center
                      , ST_BOUNDINGBOX(boundary) AS bbox
                      -- switch from az to standard angle and rotate 30 degrees counter-clockwise
                      , deg_to_rad(90 - azimuth - 30) AS alpha # alpha is rotation angle of hexbin grid
      FROM            buildings
      )
      
      , local_grid_params AS(
      -- parameters for hexbin grid, taylored to each building
      SELECT          building_id
                      , ST_X(center) AS X0
                      , ST_Y(center) AS Y0
                      , alpha
      -- Commented out by Luke Murphy 28/6/2024:
                      -- , ST_DISTANCE(ST_GEOGPOINT(bbox.xmax, bbox.ymin), ST_GEOGPOINT(bbox.xmin, bbox.ymin)) AS width
                      -- , ST_DISTANCE(ST_GEOGPOINT(bbox.xmax, bbox.ymin), ST_GEOGPOINT(bbox.xmax, bbox.ymax)) AS height
      -- Added by Luke Murphy 28/6/2024 to replace commented code above:
      -- Using the max of width or height to create a broader, square bounding box to contain unusual-shaped buildings (especially long/thin)
      -- Issue that some long/thin buildings are not fully covered by grid once rotated
      -- Also increasing the grid area by 5% to endeavour to cover all corners
                      , GREATEST(
                        ST_DISTANCE(ST_GEOGPOINT(bbox.xmax, bbox.ymin), ST_GEOGPOINT(bbox.xmin, bbox.ymin)),
                        ST_DISTANCE(ST_GEOGPOINT(bbox.xmax, bbox.ymin), ST_GEOGPOINT(bbox.xmax, bbox.ymax))
                        ) * 1.05
                        AS dist
      FROM            bboxes
      )
      
      , primary_grid AS (
      -- primary grid for hexbin centers
      SELECT        building_id
                    , X
                    , Y
                    , X0
                    , Y0
                    , alpha
      FROM          local_grid_params
      -- Commented out by Luke Murphy 28/6/2024:
                    -- , UNNEST(GENERATE_ARRAY(-CEIL(width / 2 / step_x + 2) * step_x, CEIL(width / 2 / step_x + 2) * step_x, step_x)) AS X
                    -- , UNNEST(GENERATE_ARRAY(-CEIL(height / 2 / step_y + 2) * step_y, CEIL(height / 2 / step_y + 2) * step_y, step_y)) AS Y
      -- Added by Luke Murphy 28/6/2024 to replace commented code above:
      -- X and Y now calculated using the new dist measure created above
                    , UNNEST(GENERATE_ARRAY(-CEIL(dist / 2 / step_x + 2) * step_x, CEIL(dist / 2 / step_x + 2) * step_x, step_x)) AS X
                    , UNNEST(GENERATE_ARRAY(-CEIL(dist / 2 / step_y + 2) * step_y, CEIL(dist / 2 / step_y + 2) * step_y, step_y)) AS Y
      )
      
      , full_grid AS (
      -- full grid of hexbin centers
      SELECT          GENERATE_UUID() AS uid
                      , building_id
                      , X
                      , Y
                      , X0
                      , Y0
                      , alpha
      FROM            primary_grid
      UNION ALL
      SELECT          GENERATE_UUID() AS uid
                      , building_id
                      , X + offset_x AS X
                      , Y + offset_y AS Y
                      , X0
                      , Y0
                      , alpha
      FROM            primary_grid
      )
      
      , offsets AS (
      -- maps hexbin ring nodes as offset from hexbin centre
      SELECT  1 AS x1, 0 AS y1
      UNION ALL SELECT 0.5, 1
      UNION ALL SELECT -0.5, 1
      UNION ALL SELECT -1, 0
      UNION ALL SELECT -0.5, -1
      UNION ALL SELECT 0.5, -1
      UNION ALL SELECT 1, 0
      )
      
      , nodes AS (
      -- full grid of final hexbins vertices in (X,Y) coordinates
      SELECT          uid
                      , building_id
                      , alpha
                      , ROUND(X + R * x1, 8) AS X
                      -- Y offset are in inradius...
                      , ROUND(Y + SQRT(3) / 2 * R * y1, 8) AS Y
                      , X0
                      , Y0
                      , earth_radius(Y0) Re
      FROM            full_grid AS fg
                      , offsets
      )
      
      , rotated AS (
      -- rotated hexbins vertices in (X,Y) coordinates
      SELECT        uid
                    , building_id
                    , X*COS(alpha) - Y*SIN(alpha) AS X
                    , X*SIN(alpha) + Y*COS(alpha) AS Y
                    , X0
                    , Y0
                    -- transform factors back to (lat,lon)
                    -- Rx is adjusted for radius of small circle
                    , rad_to_deg(1 / Re / COS(deg_to_rad(Y0))) AS Rx  # affine transform ratio between X and longitude
                    , rad_to_deg(1 / Re) AS Ry                        # affine transform ratio between Y and latitude
      FROM          nodes
      )
      
      , transformed AS (
      -- rotated hexbins transformed to (lat,lon) coordinates
      SELECT        uid
                    , building_id
                    , X * Rx + X0 AS X
                    , Y * Ry + Y0 AS Y
      FROM          rotated
      )
      
      , hexbins AS (
      -- hexbin polygons
      SELECT          building_id
                      , ST_CONVEXHULL(ST_UNION_AGG(ST_GEOGPOINT(X, Y))) AS geometry
      FROM            transformed
      GROUP BY        building_id, uid
      )
      
      , merged_hexbins AS (
      -- get whole building polygon geometries back
      SELECT          bldg.*
                      , hb.geometry
      FROM            hexbins AS hb
      JOIN            buildings AS bldg
                      ON hb.building_id = bldg.building_id
      )
      
      , intersections AS (
      -- clip hexbins to building outlines
      SELECT          building_id
                      , ST_INTERSECTION(boundary, geometry) AS geometry
      FROM            merged_hexbins
      )
      
      , split_buildings AS(
      -- finalized hexbin table
      SELECT          building_id
                      , geometry AS boundary
                      , ST_BUFFER(geometry, hb_buffer) AS boundary_buffered
                      , ST_AREA(geometry) AS area_m2
      FROM            intersections
      WHERE           1 = 1
      AND             NOT ST_ISEMPTY(geometry)
      )
    
    -- UNION ALL back to full set of buildings (incl. buildings that weren't split in XY)
    SELECT          GENERATE_UUID() AS uid
                    ,building_id
                    ,boundary
                    ,area_m2
                    , 1 AS is_XYsplit -- 1 IS INTERNAL BUILDING HEXBIN
                    --, 1 AS display_flag
    FROM            split_buildings AS sb
    WHERE           area_m2 > 0
    );
    
    SELECT * FROM `da-datascience-pt.LM_sandpit.Temp_LG_building_grid`;
    """
    
    # retrieve Big Query SQL as dataframe
    building_xy_grid = client.query(SQL_building_grid).to_dataframe()
    
    ###############################################################################
    ###############################################################################
    ###############################################################################
    
    
    ####### RETRIEVE MEASUREMENTS #######
    
    ## Retrieve data points for building
    SQL_building_dps = """
    SELECT  dps.Location_Altitude
            ,dps.Device_SIMServiceProviderBrandName
            ,dps.QOS_RSSNR
            ,dps.QOS_RSRP
            ,dps.stratum_no
            ,grid.uid
    FROM    `denseware-prod-1-d7b8.d_denseware_lead_generation_us_prod.t_leadgen-subset-xyz_v1` as dps,
            `da-datascience-pt.LM_sandpit.Temp_LG_building_grid` as grid
    WHERE   dps.building_id = """ + str(building_id) + """
    AND     dps.Device_SIMServiceProviderBrandName IN ('AT&T', 'T-Mobile', 'Verizon')
    AND     dps.Connection_Category = '4G'
    AND     dps.Connection_Technology = 'LTE'
    AND         st_intersects(st_geogpoint(dps.Location_Longitude, dps.Location_Latitude), grid.boundary)
    """
    
    # retrieve Big Query SQL as dataframe
    building_dps = client.query(SQL_building_dps).to_dataframe()
    
    
    #-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#
    
    
    ### TOTAL MEASUREMENTS PER XY BIN
    
    # Filter data points to main carriers
    Tot_meas_XY = building_dps[building_dps['Device_SIMServiceProviderBrandName'].isin(['AT&T', 'T-Mobile', 'Verizon'])]
    
    # Anonimise carrier labels
    conditions = [
        (Tot_meas_XY['Device_SIMServiceProviderBrandName'].str.contains("AT&T")),
        (Tot_meas_XY['Device_SIMServiceProviderBrandName'].str.contains("T-Mobile")),
        (Tot_meas_XY['Device_SIMServiceProviderBrandName'].str.contains("Verizon"))]
    
    choices = [
        'Carrier 1',
        'Carrier 2',
        'Carrier 3',
        ]
    
    # apply the lists to df
    Tot_meas_XY['carrier'] = np.select(conditions, choices, default='Other')

    # convert to floats
    Tot_meas_XY['QOS_RSSNR'] = Tot_meas_XY['QOS_RSSNR'].astype(float)    
    Tot_meas_XY['QOS_RSRP'] = Tot_meas_XY['QOS_RSSNR'].astype(float)    
    
    # Count measurements per xy bin
    xy_meas = Tot_meas_XY.groupby(['uid']).\
        agg(Measurements = ('uid', 'count')
            ).reset_index()
            
    # merge on xy bin geometry
    xy_meas = pd.merge(xy_meas[[
        'uid',
        'Measurements'
        ]],
        building_xy_grid[['uid','boundary']], on='uid', how='right').fillna(-999)
    
    # Convert the pandas DataFrame into a GeoDataFrame. 
    xy_meas['geom'] = xy_meas['boundary'].apply(wkt.loads)
    xy_meas = gpd.GeoDataFrame(xy_meas, geometry='geom').set_crs('EPSG:4326')
    
    # Apply DW colour scheme to measurements per xy bin
    xy_meas['Measurement_colour'] = pd.cut(xy_meas['Measurements'], bins=meas_bins, labels=meas_labels)
    
    # Produce graphic
    fig, axJ = plt.subplots(1, figsize=(14, 8),dpi=80)
    plt.axis('off')
    xy_meas.plot(color=xy_meas['Measurement_colour'], alpha = 1, linewidth = 0.5, edgecolor='black', legend=True, ax=axJ)
    
    fig.savefig(folder + 'Images/' + str(building_id) +'_Measurements_xy.png', transparent=True)
    
    
    #-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#
    
    
    ## PRODUCE MAP OF BUILDING LOCATION
    def image_spoof(self, tile): # this function pretends not to be a Python script
        url = self._image_url(tile) # get the url of the street map API
        req = Request(url) # start request
        req.add_header('User-agent','Anaconda 3') # add user agent to request
        fh = urlopen(req) 
        im_data = io.BytesIO(fh.read()) # get image
        fh.close() # close url
        img = Image.open(im_data) # open image with PIL
        img = img.convert(self.desired_tile_form) # set image format
        return img, self.tileextent(tile), 'lower' # reformat for cartopy
    
    
    # Zoomed out map
    cimgt.OSM.get_image = image_spoof # reformat web request for street map spoofing
    osm_img = cimgt.QuadtreeTiles() # spoofed, downloaded street map
    
    fig = plt.figure(figsize=(12,9)) # open matplotlib figure
    ax1 = plt.axes(projection=osm_img.crs) # project using coordinate reference system (CRS) of street map
    
    # Create coordinate references for centre of map using building
    x = xy_meas.dissolve().explode().centroid.x.item()
    y = xy_meas.dissolve().explode().centroid.y.item()
    
    center_pt = [y, x] # lat/lon using building of interest
    zoom = 0.3 # for zooming out of center point
    extent = [center_pt[1]-(zoom*2.0),center_pt[1]+(zoom*2.0),center_pt[0]-zoom,center_pt[0]+zoom] # adjust to zoom
    ax1.set_extent(extent) # set extents
    
    scale = np.ceil(-np.sqrt(2)*np.log(np.divide(zoom,350.0))) # empirical solve for scale based on zoom
    scale = (scale<20) and scale or 19 # scale cannot be larger than 19
    
    plt.text(x, y, '*',
             horizontalalignment='center',
             transform=ccrs.Geodetic(),
             fontsize=60,
             color = 'red')
    
    ax1.add_image(osm_img, int(scale)) # add OSM with zoom specification
    fig.savefig(folder + 'Images/' + str(building_id) +'_location_map_zoomed_out.png')
    
    
    # Zoomed in map
    cimgt.OSM.get_image = image_spoof # reformat web request for street map spoofing
    osm_img = cimgt.QuadtreeTiles() # spoofed, downloaded street map
    
    fig = plt.figure(figsize=(12,9)) # open matplotlib figure
    ax1 = plt.axes(projection=osm_img.crs) # project using coordinate reference system (CRS) of street map
    
    # Create coordinate references for centre of map using building
    x = xy_meas.dissolve().explode().centroid.x.item()
    y = xy_meas.dissolve().explode().centroid.y.item()
    
    center_pt = [y, x] # lat/lon using building of interest
    zoom = 0.003 # for zooming out of center point
    extent = [center_pt[1]-(zoom*2.0),center_pt[1]+(zoom*2.0),center_pt[0]-zoom,center_pt[0]+zoom] # adjust to zoom
    ax1.set_extent(extent) # set extents
    
    scale = np.ceil(-np.sqrt(2)*np.log(np.divide(zoom,350.0))) # empirical solve for scale based on zoom
    scale = (scale<20) and scale or 19 # scale cannot be larger than 19
    
    plt.text(x, y, '*',
             horizontalalignment='center',
             transform=ccrs.Geodetic(),
             fontsize=60,
             color = 'red')
    
    ax1.add_image(osm_img, int(scale)) # add OSM with zoom specification
    fig.savefig(folder + 'Images/' + str(building_id) +'_location_map_zoomed_in.png')
    
    
    #-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#
    
    
    ### RF FOR ENTIRE BUILDING
    
    
    ## Mean and 25th percentile for entire building per carrier
    # Count measurements per xy bin
    Tot_RF = Tot_meas_XY.groupby(['carrier']).\
        agg(rsrpMean = ('QOS_RSRP', 'mean'),
            rsrp25Pc = ('QOS_RSRP', q25),
            rssnrMean = ('QOS_RSSNR', 'mean'),
            rssnr25Pc = ('QOS_RSSNR', q25)
            ).reset_index()
    
    # Apply dw RF colour scheme
    Tot_RF['rsrpMean_colour'] = pd.cut(Tot_RF['rsrpMean'], bins=rsrp_bins, labels=rsrp_labels)
    Tot_RF['rsrp25Pc_colour'] = pd.cut(Tot_RF['rsrp25Pc'], bins=rsrp_bins, labels=rsrp_labels)
    Tot_RF['rssnrMean_colour'] = pd.cut(Tot_RF['rssnrMean'], bins=rssnr_bins, labels=rssnr_labels)
    Tot_RF['rssnr25Pc_colour'] = pd.cut(Tot_RF['rssnr25Pc'], bins=rssnr_bins, labels=rssnr_labels)

    # Remove rows containing null Mean RSSNR (as proxy denoting no analysis is avaialble)
    Tot_RF = Tot_RF[Tot_RF['rssnrMean'].notna()]
    
    # Create bar charts
    def building_RF_chart(attribute, clr, title, maxval, minval, fiddler): 
        fig = plt.figure()
        ax = fig.add_subplot(111)
        ax.bar(Tot_RF['carrier'], Tot_RF[attribute]+fiddler, bottom=-fiddler, width=0.8, color=Tot_RF[clr])
        ax.set(ylim=(maxval, minval))
        ax.invert_yaxis()
        plt.title(title, fontsize=16)
        plt.ylabel("dBm", fontsize=12)
        labels = np.round(Tot_RF[attribute],1)
        ax.bar_label(ax.containers[0], labels=labels, fontsize=10, label_type='edge')
        sns.despine(left=True)
        fig.savefig(folder + 'Images/' + str(building_id) + ' ' + title  + '.png', transparent=True)
        return()
    
    # RSRP Mean
    building_RF_chart('rsrpMean', 'rsrpMean_colour', 'Mean Coverage', -60, -140, 140)
    # RSRP 25th Percentile
    building_RF_chart('rsrp25Pc', 'rsrp25Pc_colour', 'Coverage 25th Percentile', -60, -140, 140)
    # RSSNR Mean
    building_RF_chart('rssnrMean', 'rssnrMean_colour', 'Mean Capacity', 30, -20, 20)
    # RSSNR 25th Percentile
    building_RF_chart('rssnr25Pc', 'rssnr25Pc_colour', 'Capacity 25th Percentile', 30, -20, 20)
    
    
    #-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#
    
    
    ### PER FLOOR ANALYSIS
    
    # Need to apply within python to use loops on the unknown volume of floors
     
    # Function for applying floor to dps using altitude
    # Floor height of 4.0492 is taken from regression model:
        # https://denseair.sharepoint.com/:p:/s/Denseware-Devs2/Ee8Wk2Zy1fdCv2fIi97bvkYBj3q2fiT7qiJH6IgTL5fdZA?e=JcBvzE
    def assign_floor (altitude):
        # calculate theoretical floor of altitude location relative to denseWare 10% level
        assigned_floor = math.floor((altitude - low_alt) / 4.0492)
        # ensure that lowest floor is 0, and highest reflect number of stories
        if assigned_floor < 0:
            assigned_floor = 0
        elif assigned_floor >= stories:
            assigned_floor = (stories - 1)
        else:
            assigned_floor = assigned_floor
        return(assigned_floor)
    
    
    # Subset data points to a df containing only present altitude values
    alt_dps = Tot_meas_XY.dropna(subset=['Location_Altitude'])
    
    alt_dps['est_floor'] = alt_dps.apply(lambda row: assign_floor(row['Location_Altitude']),axis=1)
    
    # Convert to account for US floor numbering
    alt_dps['est_floor'] = alt_dps['est_floor'] + 1
    
    alt_dps['est_floor'].value_counts()
    
    # Summarise to floor and MNO level
    floor_summary = alt_dps.groupby(['est_floor','carrier']).\
        agg(rsrpMean = ('QOS_RSRP', 'mean'),
            rsrp25Pc = ('QOS_RSRP', q25),
            rssnrMean = ('QOS_RSSNR', 'mean'),
            rssnr25Pc = ('QOS_RSSNR', q25),
            Measurements = ('est_floor', 'count')
            ).reset_index()
        
    # Set NaN to values that will yield no bar on subsequent charts
    floor_summary['rsrpMean'] = floor_summary['rsrpMean'].fillna(-140)
    floor_summary['rsrp25Pc'] = floor_summary['rsrp25Pc'].fillna(-140)
    floor_summary['rssnrMean'] = floor_summary['rssnrMean'].fillna(-20)
    floor_summary['rssnr25Pc'] = floor_summary['rssnr25Pc'].fillna(-20)    
    
    # Apply dw RF colour scheme
    floor_summary['rsrpMean_colour'] = pd.cut(floor_summary['rsrpMean'], bins=rsrp_bins, labels=rsrp_labels)
    floor_summary['rsrp25Pc_colour'] = pd.cut(floor_summary['rsrp25Pc'], bins=rsrp_bins, labels=rsrp_labels)
    floor_summary['rssnrMean_colour'] = pd.cut(floor_summary['rssnrMean'], bins=rssnr_bins, labels=rssnr_labels)
    floor_summary['rssnr25Pc_colour'] = pd.cut(floor_summary['rssnr25Pc'], bins=rssnr_bins, labels=rssnr_labels)

    
    # Apply marketing/experience RF colour scheme
# =============================================================================
#     floor_summary['rsrpMean_colour'] = pd.cut(floor_summary['rsrpMean'], bins=mkt_rsrp_bins, labels=mkt_rsrp_labels)
#     floor_summary['rsrp25Pc_colour'] = pd.cut(floor_summary['rsrp25Pc'], bins=mkt_rsrp_bins, labels=mkt_rsrp_labels)
#     floor_summary['rssnrMean_colour'] = pd.cut(floor_summary['rssnrMean'], bins=mkt_rssnr_bins, labels=mkt_rssnr_labels)
#     floor_summary['rssnr25Pc_colour'] = pd.cut(floor_summary['rssnr25Pc'], bins=mkt_rssnr_bins, labels=mkt_rssnr_labels)
#     
# =============================================================================
    
    # Where building > xx floors, floors will be selected by weight of measurements
    if stories < max_flrs:
        flr_meas = floor_summary.groupby(['est_floor']).\
            agg(Measurements = ('Measurements', 'sum')).reset_index()   
        flr_meas = flr_meas.sort_values(by='Measurements', ascending=False)
        flr_meas = flr_meas.head(max_flrs)
        flr_list = flr_meas['est_floor'].tolist()
        flr_list.sort()
        restr_floor_summary = floor_summary[floor_summary['est_floor'].isin(flr_list)]
        restr_floor_summary['Floor_note'] = 'Floor ' + restr_floor_summary['est_floor'].astype(str)
    else:
        flr_meas = floor_summary.groupby(['est_floor']).\
            agg(Measurements = ('Measurements', 'sum')).reset_index()   
        flr_meas = flr_meas.sort_values(by='Measurements', ascending=False)
        flr_meas = flr_meas.head(max_flrs)
        flr_list = flr_meas['est_floor'].tolist()
        flr_list.sort()
        restr_floor_summary = floor_summary[floor_summary['est_floor'].isin(flr_list)]
        restr_floor_summary['Floor_note'] = 'Floor ' + restr_floor_summary['est_floor'].astype(str)
    
    
    # Create measurement bar charts per floor
    def floor_meas_chart(MNO, attribute, title, maxval, minval): 
        fig = plt.figure()
        ax = fig.add_subplot(111)
        sample = restr_floor_summary[restr_floor_summary['carrier'].isin([MNO])]
        ax.barh(sample['Floor_note'], sample[attribute], color='lightblue', edgecolor='black', height = 1.0)
        y_pos = np.arange(len(sample['Floor_note']))
        ax.set_yticks(y_pos, labels = sample['Floor_note'], fontsize=12)
        ax.set(xlim=(maxval, minval))
        ax.invert_xaxis()
        plt.xlabel(title, fontsize=12)
        sns.despine(left=True)
        fig.savefig(folder + 'Images/' + str(building_id) + ' ' + title  + ' floor analysis for ' + MNO + '.png', transparent=True)
        return()
    
    # Measurements
    floor_meas_chart('Carrier 1', 'Measurements', 'Total measurements', restr_floor_summary['Measurements'].max()+500, 0)
    floor_meas_chart('Carrier 2', 'Measurements', 'Total measurements', restr_floor_summary['Measurements'].max()+500, 0)
    floor_meas_chart('Carrier 3', 'Measurements', 'Total measurements', restr_floor_summary['Measurements'].max()+500, 0)
    
    
    
    # Create RF bar charts for floors
    def floor_RF_chart(MNO, attribute, clr, title, maxval, minval, fiddler): 
        fig = plt.figure()
        ax = fig.add_subplot(111)
        sample = restr_floor_summary[restr_floor_summary['carrier'].isin([MNO])]
        ax.barh(sample['Floor_note'], sample[attribute] + fiddler, left = -fiddler, color=sample[clr], edgecolor='black', height = 1.0)
        y_pos = np.arange(len(sample['Floor_note']))
        ax.set_yticks(y_pos, labels = sample['Floor_note'], fontsize=12)
        ax.set(xlim=(maxval, minval))
        ax.invert_xaxis()
        plt.xlabel(title, fontsize=12)
        sns.despine(left=True)
        fig.savefig(folder + 'Images/' + str(building_id) + ' ' + title  + ' floor analysis for ' + MNO + '.png', transparent=True)
        return()
    
    # RSRP 25th
    floor_RF_chart('Carrier 1', 'rsrp25Pc', 'rsrp25Pc_colour', 'Coverage 25th Percentile', -60, -140, 140)
    floor_RF_chart('Carrier 2', 'rsrp25Pc', 'rsrp25Pc_colour', 'Coverage 25th Percentile', -60, -140, 140)
    floor_RF_chart('Carrier 3', 'rsrp25Pc', 'rsrp25Pc_colour', 'Coverage 25th Percentile', -60, -140, 140)
    
    # RSSNR 25th
    floor_RF_chart('Carrier 1', 'rssnr25Pc', 'rssnr25Pc_colour', 'Capacity 25th Percentile', 30, -20, 20)
    floor_RF_chart('Carrier 2', 'rssnr25Pc', 'rssnr25Pc_colour', 'Capacity 25th Percentile', 30, -20, 20)
    floor_RF_chart('Carrier 3', 'rssnr25Pc', 'rssnr25Pc_colour', 'Capacity 25th Percentile', 30, -20, 20)
    
    
    ###############################################################################
    ###############################################################################
    ###############################################################################
    
    
    ####### XY ANALYSIS PER FLOOR #######
    
    # Use data that has valid altitude readings
    floor_xy_summary = alt_dps.groupby(['est_floor','carrier', 'uid']).\
        agg(rsrpMean = ('QOS_RSRP', 'mean'),
            rsrp25Pc = ('QOS_RSRP', q25),
            rssnrMean = ('QOS_RSSNR', 'mean'),
            rssnr25Pc = ('QOS_RSSNR', q25)
            ).reset_index()
        
    # limit to selected floors
    floor_xy_summary = floor_xy_summary[floor_xy_summary['est_floor'].isin(flr_list)]
        
    
    # Function for XY graphic
    def xy_rf_per_floor(MNO, clrmap, title, flr): 
        fig, axJ = plt.subplots(1, figsize=(14, 8),dpi=80)
        plt.axis('off')
        dat = floor_xy_summary[(
            (floor_xy_summary['carrier'].isin([MNO])) & 
            (floor_xy_summary['est_floor'] == flr)
             )]
        
        dat = pd.merge(
            building_xy_grid[['uid','boundary']],
            dat[['uid',
                'rsrpMean',
                'rsrp25Pc',
                'rssnrMean',
                'rssnr25Pc'
                ]],
            on='uid', how='left')
        
        dat['rsrpMean'] = dat['rsrpMean'].fillna(999)
        dat['rsrp25Pc'] = dat['rsrp25Pc'].fillna(999)
        dat['rssnrMean'] = dat['rssnrMean'].fillna(999)
        dat['rssnr25Pc'] = dat['rssnr25Pc'].fillna(999)
    
        dat['geom'] = dat['boundary'].apply(wkt.loads)
        dat = gpd.GeoDataFrame(dat, geometry='geom').set_crs('EPSG:4326')
        
        dat['rsrpMean_colour'] = pd.cut(dat['rsrpMean'], bins=rsrp_bins, labels=rsrp_labels)
        dat['rsrp25Pc_colour'] = pd.cut(dat['rsrp25Pc'], bins=rsrp_bins, labels=rsrp_labels)
        dat['rssnrMean_colour'] = pd.cut(dat['rssnrMean'], bins=rssnr_bins, labels=rssnr_labels)
        dat['rssnr25Pc_colour'] = pd.cut(dat['rssnr25Pc'], bins=rssnr_bins, labels=rssnr_labels)
        
        dat.plot(color=dat[clrmap], alpha = 1, linewidth = 0.5, edgecolor='black', legend=True, ax=axJ)
        fig.savefig(folder + 'Images/' + str(building_id) + ' ' + title  + ' xy floor ' + str(flr) + ' ' + MNO + '.png', transparent=True)
        return()  
    
    # Iterate through number of floors to produce appropriate grpahics
    
    for i, flr in enumerate(flr_list):  
        # RSRP Mean
        xy_rf_per_floor('Carrier 1', 'rsrpMean_colour', 'Mean Coverage', flr)
        xy_rf_per_floor('Carrier 2', 'rsrpMean_colour', 'Mean Coverage', flr)
        xy_rf_per_floor('Carrier 3', 'rsrpMean_colour', 'Mean Coverage', flr)
    
        # RSRP 25th
        xy_rf_per_floor('Carrier 1', 'rsrp25Pc_colour', 'Coverage 25th Percentile', flr)
        xy_rf_per_floor('Carrier 2', 'rsrp25Pc_colour', 'Coverage 25th Percentile', flr)
        xy_rf_per_floor('Carrier 3', 'rsrp25Pc_colour', 'Coverage 25th Percentile', flr)
    
        # RSSNR Mean
        xy_rf_per_floor('Carrier 1', 'rssnrMean_colour', 'Mean Capacity', flr)
        xy_rf_per_floor('Carrier 2', 'rssnrMean_colour', 'Mean Capacity', flr)
        xy_rf_per_floor('Carrier 3', 'rssnrMean_colour', 'Mean Capacity', flr)
    
        # RSSNR 25th
        xy_rf_per_floor('Carrier 1', 'rssnr25Pc_colour', 'Capacity 25th Percentile', flr)
        xy_rf_per_floor('Carrier 2', 'rssnr25Pc_colour', 'Capacity 25th Percentile', flr)
        xy_rf_per_floor('Carrier 3', 'rssnr25Pc_colour', 'Capacity 25th Percentile', flr)
    
    
    
    ###############################################################################
    ###############################################################################
    ###############################################################################
    
    
    ####### CREATE AND SAVE POWERPOINT #######
    
    #### Title slide
    prs = Presentation(folder + 'Templates/' + template)
    
    #### Title slide
    # prs.slides command used to identify the slide that you wish to alter
    slide = prs.slides[0]
    title = slide.shapes.title
    title.text = "denseWare building analysis"
    subtitle = slide.placeholders[1]
    subtitle.text = address + ', ' + city + ', ' + state + ', ' + zip_code
    
    # END OF SLIDE
    
    #-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#
    
    #### Overview slide
    
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    
    # Add title
    slide.shapes.title.text = "Overview"
    
    ## Create a subtitle
    txBox = slide.shapes.add_textbox(Cm(1), Cm(1.5), width = Cm(10), height = Cm(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = address + ', ' + zip_code
    p.font.name = 'Avenir Next LT Pro'
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(38, 76, 143)
     
    ## Comments on slide
    txBox = slide.shapes.add_textbox(Cm(1), Cm(8.2), width = Cm(16), height = Cm(10))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    
    # first 'bullet'
    p.text = "•" + " The property is situated in " + city + ", " + state + "."
    p.font.name = 'Avenir Next LT Pro'
    p.font.size = Pt(14)
    
    # add a space
    p = tf.add_paragraph()
    p.text = ''
    
    # second 'bullet' etc etc
    p = tf.add_paragraph()
    p.text = "•" + " CoStar property report details " + str(stories) \
        + " stories, with a total RBA (rentable building area) of " + str(f'{rba:,}') \
            + " square feet."
    p.font.name = 'Avenir Next LT Pro'
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = ''
    
    p = tf.add_paragraph()
    p.text = "•" + " The property was built, or most recently renovated, in " + str(year) +"."
    p.font.name = 'Avenir Next LT Pro'
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = ''
    
    p = tf.add_paragraph()
    p.text = "•" + " Data used for analysis includes "+f'{len(Tot_meas_XY.index):,}'+" crowdsource data measurements."
    p.font.name = 'Avenir Next LT Pro'
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = ''
    
    p = tf.add_paragraph()
    p.text = "•" + " INSERT YOUR OBSERVATIONS HERE."
    p.font.name = 'Avenir Next LT Pro'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.italic = True
    
    # call hexbin size image
    slide.shapes.add_picture(hxbin_size_image, Cm(21), Cm(18), height = Cm(0.8))    
    txBox = slide.shapes.add_textbox(Cm(22), Cm(17), width = Cm(5.5), height = Cm(2))
    tf2 = txBox.text_frame
    tf2.word_wrap = True
    p2 = tf2.add_paragraph()
    p2.text = "Network data are analyzed inside hexbins of " + str(hexbin_size) + " sq ft, measuring " + str(hexbin_diag) + " ft across long diagonal."
    p2.font.name = 'Avenir Next LT Pro'
    p2.font.size = Pt(9)
    p2.font.color.rgb = RGBColor(255, 255, 255)
    
               
    # apply zoomed IN location map image
    Img_meas = folder + 'Images/' + str(building_id) +'_location_map_zoomed_in.png'
    slide.shapes.add_picture(Img_meas, Cm(21.5), Cm(0), height = Cm(9))
    
    # apply zoomed OUT location map image
    Img_meas = folder + 'Images/' + str(building_id) +'_location_map_zoomed_out.png'
    slide.shapes.add_picture(Img_meas, Cm(11), Cm(0), height = Cm(9))
    
    # apply xy measurements image
    Img_meas = folder + 'Images/' + str(building_id) +'_Measurements_xy.png'
    slide.shapes.add_picture(Img_meas, Cm(13), Cm(8), height = Cm(11))
    
    # call measurement scale image
    slide.shapes.add_picture(meas_scale_image, Cm(29), Cm(11), height = Cm(6))        
    
    
    # END OF SLIDE
    
    #-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#
    
    #### Network performance slide
    
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    
    # Add title
    slide.shapes.title.text = 'Network performance - property overview'
    
    # Set Ref df from building summary to call in logic statements
    ref_df = Tot_RF.set_index('carrier')
    
    ## Create a subtitle
    txBox = slide.shapes.add_textbox(Cm(1), Cm(1.5), width = Cm(31), height = Cm(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "Summary network KPI statistics by carrier using crowdsourced network measurements, with a focus on LTE."
    p.font.name = 'Avenir Next LT Pro'
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(38, 76, 143)
     
    ## Comments on slide
    txBox = slide.shapes.add_textbox(Cm(1), Cm(3), width = Cm(15.5), height = Cm(12))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    
    # first 'bullet'
    p.text = "•" + " RSRP (signal power) is a measure of coverage, indicating how far RF waves propagate. Higher is better, with levels lower than -108dBm indicative of poor performance, and equivalent ~1 ‘bar’ on a mobile device signal strength gauge."
    p.font.name = 'Avenir Next LT Pro'
    p.font.size = Pt(14)
    
    # add a space
    p = tf.add_paragraph()
    p.text = ''
    
    # second 'bullet' etc etc
    p = tf.add_paragraph()
    p.text = "•" + " RSSNR (signal-to-noise) is a measure of capacity indicating how much of the RF waves can be used to transmit data. Higher is better, with values lower than 0dB indicative of poor performance.​"
    p.font.name = 'Avenir Next LT Pro'
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = ''
    
    p = tf.add_paragraph()
    p.text = "•" + " The 25th Percentile represents the experience of 1 in 4 users and tends to represent user experience during times of peak network load."
    p.font.name = 'Avenir Next LT Pro'
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = ''
    
    # Number of rows in table to analyse
    carriers = len(Tot_RF.index)
 
    p = tf.add_paragraph()
    if carriers >= 1:
        if ref_df.at['Carrier 1', 'rsrp25Pc'] <= -108 and ref_df.at['Carrier 1', 'rssnr25Pc'] <= 0:
            p.text = "•" + " Carrier 1 demonstrates both coverage and capacity issues at peak times."
        elif ref_df.at['Carrier 1', 'rsrp25Pc'] <= -108 and ref_df.at['Carrier 1', 'rssnr25Pc'] > 0:
            p.text = "•" + " Carrier 1 demonstrates coverage issues at peak times."
        elif ref_df.at['Carrier 1', 'rsrp25Pc'] > -108 and ref_df.at['Carrier 1', 'rssnr25Pc'] <= 0:
            p.text = "•" + " Carrier 1 demonstrates capacity issues at peak times."       
        else:
            p.text = "•" + " Carrier 1 has no apparent coverage or capacity issues at peak times."
    else:
        p.text = "•" + " No data available for carrier 1."  
    p.font.name = 'Avenir Next LT Pro'
    p.font.size = Pt(14)
    p.font.bold = True 
    p = tf.add_paragraph()
    p.text = ''
    
    
    p = tf.add_paragraph()
    if carriers >= 2:
        if ref_df.at['Carrier 2', 'rsrp25Pc'] <= -108 and ref_df.at['Carrier 2', 'rssnr25Pc'] <= 0:
            p.text = "•" + " Carrier 2 demonstrates both coverage and capacity issues at peak times."
        elif ref_df.at['Carrier 2', 'rsrp25Pc'] <= -108 and ref_df.at['Carrier 2', 'rssnr25Pc'] > 0:
            p.text = "•" + " Carrier 2 demonstrates coverage issues at peak times."
        elif ref_df.at['Carrier 2', 'rsrp25Pc'] > -108 and ref_df.at['Carrier 2', 'rssnr25Pc'] <= 0:
            p.text = "•" + " Carrier 2 demonstrates capacity issues at peak times."       
        else:
            p.text = "•" + " Carrier 2 has no apparent coverage or capacity issues at peak times."
    else:
        p.text = "•" + " No data available for carrier 2."  
    p.font.name = 'Avenir Next LT Pro'
    p.font.size = Pt(14)
    p.font.bold = True 
    p = tf.add_paragraph()
    p.text = ''  
    
    
    p = tf.add_paragraph()
    if carriers >= 3:
        if ref_df.at['Carrier 3', 'rsrp25Pc'] <= -108 and ref_df.at['Carrier 3', 'rssnr25Pc'] <= 0:
            p.text = "•" + " Carrier 3 demonstrates both coverage and capacity issues at peak times."
        elif ref_df.at['Carrier 3', 'rsrp25Pc'] <= -108 and ref_df.at['Carrier 3', 'rssnr25Pc'] > 0:
            p.text = "•" + " Carrier 3 demonstrates coverage issues at peak times."
        elif ref_df.at['Carrier 3', 'rsrp25Pc'] > -108 and ref_df.at['Carrier 3', 'rssnr25Pc'] <= 0:
            p.text = "•" + " Carrier 3 demonstrates capacity issues at peak times."       
        else:
            p.text = "•" + " Carrier 3 has no apparent coverage or capacity issues at peak times."
    else:
        p.text = "•" + " No data available for carrier 3."  
    p.font.name = 'Avenir Next LT Pro'
    p.font.size = Pt(14)
    p.font.bold = True 
    p = tf.add_paragraph()
    p.text = ''
    
    
    # apply RSRP mean image
    img_rsrp_mean = folder + 'Images/' + str(building_id) + ' Mean Coverage.png'
    slide.shapes.add_picture(img_rsrp_mean, Cm(17), Cm(4), height = Cm(6))
    
    # apply RSRP 25th image
    img_rsrp_25 = folder + 'Images/' + str(building_id) + ' Coverage 25th Percentile.png'
    slide.shapes.add_picture(img_rsrp_25, Cm(25), Cm(4), height = Cm(6))
    
    # apply RSSNR mean image
    img_rsrp_mean = folder + 'Images/' + str(building_id) + ' Mean Capacity.png'
    slide.shapes.add_picture(img_rsrp_mean, Cm(17), Cm(10), height = Cm(6))
    
    # apply RSSNR 25th image
    img_rsrp_25 = folder + 'Images/' + str(building_id) + ' Capacity 25th Percentile.png'
    slide.shapes.add_picture(img_rsrp_25, Cm(25), Cm(10), height = Cm(6))
    
    # END OF SLIDE
    
    #-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#
    
    #### Overview of network performance by floor for entire building
    
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    
    # Add title
    slide.shapes.title.text = 'Network performance - floor analysis'
        
    ## Comments on slide
    txBox = slide.shapes.add_textbox(Cm(1), Cm(3), width = Cm(12), height = Cm(12))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    
    # first 'bullet'
    p.text = "•" + " Network measurements are pinned to stories to assess RF performance per property floor."
    p.font.name = 'Avenir Next LT Pro'
    p.font.size = Pt(14)
    
    # add a space
    p = tf.add_paragraph()
    p.text = ''
    
    # second 'bullet' etc etc
    p = tf.add_paragraph()
    p.text = "•" + " Floors are specified during NQT testing​."
    p.font.name = 'Avenir Next LT Pro'
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = ''
    
    p = tf.add_paragraph()
    p.text = "•" + " For crowdsourced data, floor is assigned through its altitude reading."
    p.font.name = 'Avenir Next LT Pro'
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = ''
    
    p = tf.add_paragraph()
    if stories < max_flrs:
        p.text = "•" + " Performance for each floor in the property is displayed in the charts."
    else:
        p.text = "•" + " The property has more than " + str(max_flrs) + " stories. Therefore, performance is displayed for the " + str(max_flrs) + " stories having the highest number of measurements."
    p.font.name = 'Avenir Next LT Pro'
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = ''
    
    p = tf.add_paragraph()
    p.text = "•" + " INSERT YOUR OBSERVATIONS HERE."
    p.font.name = 'Avenir Next LT Pro'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.italic = True
    
    
    ## Add labels
    txBox = slide.shapes.add_textbox(Cm(16.5), Cm(2), width = Cm(3), height = Cm(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Count of Measurements"
    p.font.name = 'Avenir Next LT Pro'
    p.font.size = Pt(12)
    p.font.underline = True
    p.font.bold = True
        
    txBox = slide.shapes.add_textbox(Cm(22), Cm(2), width = Cm(3), height = Cm(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Coverage 25th Percentile"
    p.font.name = 'Avenir Next LT Pro'
    p.font.size = Pt(12)
    p.font.underline = True
    p.font.bold = True
    
    txBox = slide.shapes.add_textbox(Cm(27.5), Cm(2), width = Cm(3), height = Cm(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Capacity 25th Percentile"
    p.font.name = 'Avenir Next LT Pro'
    p.font.size = Pt(12)
    p.font.underline = True
    p.font.bold = True
    
    txBox = slide.shapes.add_textbox(Cm(13), Cm(4), width = Cm(4), height = Cm(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "Carrier 1"
    p.font.name = 'Avenir Next LT Pro'
    p.font.size = Pt(12)
    p.font.underline = True
    p.font.bold = True
    
    txBox = slide.shapes.add_textbox(Cm(13), Cm(9), width = Cm(4), height = Cm(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "Carrier 2"
    p.font.name = 'Avenir Next LT Pro'
    p.font.size = Pt(12)
    p.font.underline = True
    p.font.bold = True 
    
    txBox = slide.shapes.add_textbox(Cm(13), Cm(14), width = Cm(4), height = Cm(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "Carrier 3"
    p.font.name = 'Avenir Next LT Pro'
    p.font.size = Pt(12)
    p.font.underline = True
    p.font.bold = True 
    p.border = True 
    
    # location of first chart (others will reference this)
    left = 16.5
    top = 4
    height = 3.5
    
    # apply images - Measurements
    Meas_c1 = folder + 'Images/' + str(building_id) + ' ' + 'Total measurements floor analysis for Carrier 1.png'
    Meas_c2 = folder + 'Images/' + str(building_id) + ' ' + 'Total measurements floor analysis for Carrier 2.png'
    Meas_c3 = folder + 'Images/' + str(building_id) + ' ' + 'Total measurements floor analysis for Carrier 3.png'
    
    slide.shapes.add_picture(Meas_c1, Cm(left), Cm(top), height = Cm(height))
    slide.shapes.add_picture(Meas_c2, Cm(left), Cm(top+5), height = Cm(height))
    slide.shapes.add_picture(Meas_c3, Cm(left), Cm(top+10), height = Cm(height))
    
    
    # apply images - RSRP 25th
    rsrp_c1 = folder + 'Images/' + str(building_id) + ' ' + 'Coverage 25th Percentile floor analysis for Carrier 1.png'
    rsrp_c2 = folder + 'Images/' + str(building_id) + ' ' + 'Coverage 25th Percentile floor analysis for Carrier 2.png'
    rsrp_c3 = folder + 'Images/' + str(building_id) + ' ' + 'Coverage 25th Percentile floor analysis for Carrier 3.png'
    
    slide.shapes.add_picture(rsrp_c1, Cm(left+5.5), Cm(top), height = Cm(height))
    slide.shapes.add_picture(rsrp_c2, Cm(left+5.5), Cm(top+5), height = Cm(height))
    slide.shapes.add_picture(rsrp_c3, Cm(left+5.5), Cm(top+10), height = Cm(height))
    
    
    # apply images - RSSNR 25th
    rssnr_c1 = folder + 'Images/' + str(building_id) + ' ' + 'Capacity 25th Percentile floor analysis for Carrier 1.png'
    rssnr_c2 = folder + 'Images/' + str(building_id) + ' ' + 'Capacity 25th Percentile floor analysis for Carrier 2.png'
    rssnr_c3 = folder + 'Images/' + str(building_id) + ' ' + 'Capacity 25th Percentile floor analysis for Carrier 3.png'
    
    slide.shapes.add_picture(rssnr_c1, Cm(left+11), Cm(top), height = Cm(height))
    slide.shapes.add_picture(rssnr_c2, Cm(left+11), Cm(top+5), height = Cm(height))
    slide.shapes.add_picture(rssnr_c3, Cm(left+11), Cm(top+10), height = Cm(height))
    
    
    # END OF SLIDE
    
    #-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#
    
    #### Network performance for each floor
    ## Creates a slide for each floor
    #slide = prs.slides[4]
    #title = slide.shapes.title
    #title.text = "Network performance - FIRST LEVEL"
    
    # location of first xy map (others will reference this)
    left = 2
    top = 2
    height = 5.5
    
    # function to apply xy images
    def add_slide(prs, layout, title, flr):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title + ' ' + str(flr)
        
        ## Add labels
        txBox = slide.shapes.add_textbox(Cm(5), Cm(1.5), width = Cm(3), height = Cm(1))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = "Mean Coverage"
        p.font.name = 'Avenir Next LT Pro'
        p.font.size = Pt(14)
        p.font.underline = True
        p.font.bold = True
        
        txBox = slide.shapes.add_textbox(Cm(11), Cm(1.5), width = Cm(3), height = Cm(1))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = "25th Percentile Coverage"
        p.font.name = 'Avenir Next LT Pro'
        p.font.size = Pt(14)
        p.font.underline = True
        p.font.bold = True
        
        txBox = slide.shapes.add_textbox(Cm(19), Cm(1.5), width = Cm(3), height = Cm(1))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = "Mean Capacity"
        p.font.name = 'Avenir Next LT Pro'
        p.font.size = Pt(14)
        p.font.underline = True
        p.font.bold = True
        
        txBox = slide.shapes.add_textbox(Cm(24), Cm(1.5), width = Cm(3), height = Cm(1))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = "25th Percentile Capacity"
        p.font.name = 'Avenir Next LT Pro'
        p.font.size = Pt(14)
        p.font.underline = True
        p.font.bold = True
        
        txBox = slide.shapes.add_textbox(Cm(1), Cm(4), width = Cm(3), height = Cm(1))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = "Carrier 1"
        p.font.name = 'Avenir Next LT Pro'
        p.font.size = Pt(14)
        p.font.underline = True
        p.font.bold = True
        
        txBox = slide.shapes.add_textbox(Cm(1), Cm(9), width = Cm(3), height = Cm(1))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = "Carrier 2"
        p.font.name = 'Avenir Next LT Pro'
        p.font.size = Pt(14)
        p.font.underline = True
        p.font.bold = True
    
        txBox = slide.shapes.add_textbox(Cm(1), Cm(14), width = Cm(3), height = Cm(1))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = "Carrier 3"
        p.font.name = 'Avenir Next LT Pro'
        p.font.size = Pt(14)
        p.font.underline = True
        p.font.bold = True
        
        ## Carrier 1
        img_rsrp_mean = folder + 'Images/' + str(building_id) + ' ' + 'Mean Coverage xy floor ' + str(flr) + ' Carrier 1.png'
        slide.shapes.add_picture(img_rsrp_mean, Cm(left), Cm(top), height = Cm(height))
        
        img_rsrp_25 = folder + 'Images/' + str(building_id) + ' ' + 'Coverage 25th Percentile xy floor ' + str(flr) + ' Carrier 1.png'
        slide.shapes.add_picture(img_rsrp_25, Cm(left+6), Cm(top), height = Cm(height))
    
        img_rssnr_mean = folder + 'Images/' + str(building_id) + ' ' + 'Mean Capacity xy floor ' + str(flr) + ' Carrier 1.png'
        slide.shapes.add_picture(img_rssnr_mean, Cm(left+13), Cm(top), height = Cm(height))
        
        img_rssnr_25 = folder + 'Images/' + str(building_id) + ' ' + 'Capacity 25th Percentile xy floor ' + str(flr) + ' Carrier 1.png'
        slide.shapes.add_picture(img_rssnr_25, Cm(left+20), Cm(top), height = Cm(height))
    
        ## Carrier 2
        img_rsrp_mean = folder + 'Images/' + str(building_id) + ' ' + 'Mean Coverage xy floor ' + str(flr) + ' Carrier 2.png'
        slide.shapes.add_picture(img_rsrp_mean, Cm(left), Cm(top+5.5), height = Cm(height))
        
        img_rsrp_25 = folder + 'Images/' + str(building_id) + ' ' + 'Coverage 25th Percentile xy floor ' + str(flr) + ' Carrier 2.png'
        slide.shapes.add_picture(img_rsrp_25, Cm(left+6), Cm(top+5.5), height = Cm(height))
    
        img_rssnr_mean = folder + 'Images/' + str(building_id) + ' ' + 'Mean Capacity xy floor ' + str(flr) + ' Carrier 2.png'
        slide.shapes.add_picture(img_rssnr_mean, Cm(left+13), Cm(top+5.5), height = Cm(height))
        
        img_rssnr_25 = folder + 'Images/' + str(building_id) + ' ' + 'Capacity 25th Percentile xy floor ' + str(flr) + ' Carrier 2.png'
        slide.shapes.add_picture(img_rssnr_25, Cm(left+20), Cm(top+5.5), height = Cm(height))    
        
        ## Carrier 3
        img_rsrp_mean = folder + 'Images/' + str(building_id) + ' ' + 'Mean Coverage xy floor ' + str(flr) + ' Carrier 3.png'
        slide.shapes.add_picture(img_rsrp_mean, Cm(left), Cm(top+11), height = Cm(height))
        
        img_rsrp_25 = folder + 'Images/' + str(building_id) + ' ' + 'Coverage 25th Percentile xy floor ' + str(flr) + ' Carrier 3.png'
        slide.shapes.add_picture(img_rsrp_25, Cm(left+6), Cm(top+11), height = Cm(height))
    
        img_rssnr_mean = folder + 'Images/' + str(building_id) + ' ' + 'Mean Capacity xy floor ' + str(flr) + ' Carrier 3.png'
        slide.shapes.add_picture(img_rssnr_mean, Cm(left+13), Cm(top+11), height = Cm(height))
        
        img_rssnr_25 = folder + 'Images/' + str(building_id) + ' ' + 'Capacity 25th Percentile xy floor ' + str(flr) + ' Carrier 3.png'
        slide.shapes.add_picture(img_rssnr_25, Cm(left+20), Cm(top+11), height = Cm(height))        
        
    # call hexbin size image
        slide.shapes.add_picture(hxbin_size_image, Cm(21), Cm(18), height = Cm(0.8))    
        txBox = slide.shapes.add_textbox(Cm(22), Cm(17), width = Cm(5.5), height = Cm(2))
        tf2 = txBox.text_frame
        tf2.word_wrap = True
        p2 = tf2.add_paragraph()
        p2.text = "Network data are analyzed inside hexbins of " + str(hexbin_size) + " sq ft, measuring " + str(hexbin_diag) + " ft across long diagonal."
        p2.font.name = 'Avenir Next LT Pro'
        p2.font.size = Pt(9)
        p2.font.color.rgb = RGBColor(255, 255, 255)
        
        # RF colour scale
        slide.shapes.add_picture(RF_scale_image, Cm(31.5), Cm(6), height = Cm(8))        
    
        return slide
    
    # choose slide layout
    slide_layout = prs.slide_layouts[2]
    
    # iterate through floors, creating a slide for each and applying xy images
    for i, flr in enumerate(flr_list):  
        add_slide(prs, slide_layout, 'Analysis for floor', flr)
    
    # END OF SLIDE(S)
    
    #-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#
    
    
    #### Output final ppt
    
    # Create datestamp for output
    now = datetime.now()
    stamp = now.strftime("%Y") + now.strftime("%m")  + now.strftime("%d") + "_" + now.strftime("%H") + "_" + now.strftime("%M") + "_" + now.strftime("%S")
    
    prs.save(folder + 'Output/' + address + ' ' \
             + city + ' ' + state + ' ' + str(building_id) + ' analysis pack ' + stamp + '.pptx')
    
        
    #-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#-#
    
    
    #### Remove image files from folder
        
    for p in Path(folder + "Images/").glob(str(building_id) + "*.png"):
        p.unlink()


##############################################################################
##############################################################################
############################# END - OF - PROGRAM #############################
##############################################################################
##############################################################################
